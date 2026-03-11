"""
Generate a professional PowerPoint presentation for the CS Ops Dashboard panel presentation.
Run this script once to create the .pptx file.

Usage (in Docker or with Python installed):
    python generate_presentation.py

Output: CS_Ops_Panel_Presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Color Palette ---
NAVY = RGBColor(15, 43, 70)       # #0F2B46
NAVY_LIGHT = RGBColor(26, 63, 92) # #1A3F5C
CYAN = RGBColor(0, 180, 216)      # #00B4D8
PURPLE = RGBColor(123, 97, 255)   # #7B61FF
GREEN = RGBColor(6, 214, 160)     # #06D6A0
AMBER = RGBColor(255, 183, 3)     # #FFB703
RED = RGBColor(239, 71, 111)      # #EF476F
WHITE = RGBColor(255, 255, 255)
DARK_TEXT = RGBColor(15, 23, 42)   # #0F172A
MUTED_TEXT = RGBColor(51, 65, 85)  # #334155
LIGHT_BG = RGBColor(241, 245, 249) # #F1F5F9
BORDER_COLOR = RGBColor(226, 232, 240)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_background(slide, color=WHITE):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_points(slide, left, top, width, height, items, font_size=16, color=DARK_TEXT, bold_prefix=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        p.space_before = Pt(3)

        if bold_prefix and ": " in item:
            prefix, rest = item.split(": ", 1)
            run1 = p.add_run()
            run1.text = prefix + ": "
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = color
            run1.font.bold = True
            run1.font.name = "Calibri"
            run2 = p.add_run()
            run2.text = rest
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.name = "Calibri"
        else:
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = "Calibri"
    return txBox


def add_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_background(slide, NAVY)
add_shape_bg(slide, Inches(0), Inches(6.8), Inches(13.333), Inches(0.06), CYAN)

add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1),
             "CS OPS INTELLIGENCE DASHBOARD", 14, CYAN, True)
add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.5),
             "Operationalizing Customer Success\nwith Data-Driven Intelligence", 40, WHITE, True)
add_text_box(slide, Inches(1.5), Inches(4.3), Inches(10), Inches(0.5),
             "Panel Presentation  |  Senior Customer Operations Analyst", 18, RGBColor(148, 163, 184))
add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
             "Sumo Logic  |  Customer Success Organization", 16, RGBColor(100, 116, 139))

add_notes(slide, """OPENING SCRIPT:
"Thank you [names] for your time today. I'm excited to be here.

When I looked at this role — Senior Customer Operations Analyst — I asked myself: What does world-class CS Operations actually look like in practice?

The answer I kept coming back to is: it's about turning customer data into action before it's too late.

So instead of just talking about what I would do, I built something to show you."
""")


# ============================================================
# SLIDE 2: Agenda
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), CYAN)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.7),
             "AGENDA", 12, CYAN, True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.8),
             "Presentation Overview", 32, NAVY, True)

agenda_items = [
    "1.  Opening & Context Setting  (5 min)",
    "2.  Executive Summary — Portfolio KPIs  (8 min)",
    "3.  Renewal Forecast & Pipeline  (7 min)",
    "4.  Churn Analysis & Risk Intelligence  (7 min)",
    "5.  Account Health Overview  (5 min)",
    "6.  EBR Preparation Workflow  (7 min)",
    "7.  PS & Training Analytics  (5 min)",
    "8.  Automated Alerts Engine  (6 min)",
    "9.  Closing: Vision & Impact  (5 min)",
    "10. Q&A  (5-10 min)",
]
add_bullet_points(slide, Inches(1.2), Inches(2.0), Inches(7), Inches(5), agenda_items, 18, DARK_TEXT)

add_notes(slide, """Walk through the agenda quickly. Mention total time is 45-60 minutes.
Say: "I've structured this around a live demo of the dashboard I built. Each section solves a real CS Ops problem."
""")


# ============================================================
# SLIDE 3: The Panel — Tailored Value
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.7),
             "YOUR PANEL", 12, PURPLE, True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.8),
             "Tailored to Each Stakeholder", 32, NAVY, True)

# Card 1 - CCO
add_shape_bg(slide, Inches(0.8), Inches(2.2), Inches(3.6), Inches(4.2), LIGHT_BG)
add_shape_bg(slide, Inches(0.8), Inches(2.2), Inches(3.6), Inches(0.08), CYAN)
add_text_box(slide, Inches(1.1), Inches(2.5), Inches(3), Inches(0.5),
             "CCO", 22, NAVY, True)
add_bullet_points(slide, Inches(1.1), Inches(3.2), Inches(3.2), Inches(3),
    ["ARR retention & growth", "NRR / GRR metrics", "Board-ready reporting", "Scalable CS operations", "Strategic churn reduction"],
    14, MUTED_TEXT)

# Card 2 - Support Ops
add_shape_bg(slide, Inches(4.9), Inches(2.2), Inches(3.6), Inches(4.2), LIGHT_BG)
add_shape_bg(slide, Inches(4.9), Inches(2.2), Inches(3.6), Inches(0.08), GREEN)
add_text_box(slide, Inches(5.2), Inches(2.5), Inches(3), Inches(0.5),
             "Dir. Global Support Ops", 18, NAVY, True)
add_bullet_points(slide, Inches(5.2), Inches(3.2), Inches(3.2), Inches(3),
    ["Escalation visibility", "Ticket trend insights", "Support-to-CS handoff", "Operational efficiency", "Reducing firefighting"],
    14, MUTED_TEXT)

# Card 3 - TAE
add_shape_bg(slide, Inches(9.0), Inches(2.2), Inches(3.6), Inches(4.2), LIGHT_BG)
add_shape_bg(slide, Inches(9.0), Inches(2.2), Inches(3.6), Inches(0.08), PURPLE)
add_text_box(slide, Inches(9.3), Inches(2.5), Inches(3), Inches(0.5),
             "Sr. Dir. TAE", 22, NAVY, True)
add_bullet_points(slide, Inches(9.3), Inches(3.2), Inches(3.2), Inches(3),
    ["Account health signals", "TAE portfolio insights", "Renewal readiness", "PS utilization", "Proactive engagement"],
    14, MUTED_TEXT)

add_notes(slide, """DO NOT show this slide to the panel — this is for YOUR reference only.
It reminds you what each person cares about so you can tailor your demo commentary in real-time.

Alternatively, you can briefly acknowledge each panelist's perspective:
"I've thought about what matters most to each of you, and I've built this dashboard with all three perspectives in mind."
""")


# ============================================================
# SLIDE 4: Context — Why I Built This
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), CYAN)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.7),
             "CONTEXT", 12, CYAN, True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.8),
             "Why I Built This Dashboard", 32, NAVY, True)

add_bullet_points(slide, Inches(1.2), Inches(2.2), Inches(10), Inches(4.5), [
    "Instead of talking about what I would do, I built something to show you",
    "",
    "Every page solves a real CS Ops operational problem:",
    "    How do we forecast renewals with confidence?",
    "    How do we catch churn signals early enough to act?",
    "    How do we prepare CSMs and TAEs for executive conversations?",
    "    How do we replace manual weekly reviews with automated intelligence?",
    "",
    "Built specifically with Sumo Logic's SaaS business model in mind",
    "250 accounts  |  27 months of trending data  |  7 operational views",
], 18, DARK_TEXT)

add_notes(slide, """KEY MESSAGE: You didn't just prepare a PowerPoint — you built a working prototype.

Say: "This isn't just a visualization exercise. Every page ends with something someone can DO — a list to export, an EBR to generate, an alert to act on."

Transition: "Let me switch to the live dashboard and walk you through it."

*** SWITCH TO LIVE DASHBOARD NOW — http://localhost:8501 ***
""")


# ============================================================
# SLIDE 5: LIVE DEMO — Executive Summary
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)
add_shape_bg(slide, Inches(0), Inches(3.5), Inches(13.333), Inches(0.06), CYAN)

add_text_box(slide, Inches(1), Inches(0.5), Inches(3), Inches(0.5),
             "LIVE DEMO", 14, CYAN, True)
add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1),
             "Executive Summary", 44, WHITE, True)
add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.5),
             "Real-time overview of Customer Success Operations KPIs", 20, RGBColor(148, 163, 184))

add_bullet_points(slide, Inches(1), Inches(4.0), Inches(5.5), Inches(3), [
    "CCO: Total ARR, Churn Rate, NRR Trend",
    "Support Ops: At-Risk count, Health Distribution",
    "TAE: ARR Waterfall, Renewal Rate Trend",
], 16, RGBColor(203, 213, 225), bold_prefix=True)

add_text_box(slide, Inches(7.5), Inches(4.2), Inches(5), Inches(2),
             '"This is what a CS Ops leader sees first thing Monday morning — a single view of everything that matters."',
             16, RGBColor(148, 163, 184))

add_notes(slide, """SWITCH TO DASHBOARD — Executive Summary page

For the CCO:
- Point to Total ARR and Active Accounts: "Portfolio health at a glance in under 2 seconds."
- Churn Rate with Target: "CS Ops should always measure against a goal, not just report a number."
- ARR Trend + NRR: "This is the story you tell the board."

For Support Ops Director:
- At-Risk Accounts: "This number directly impacts your team — these generate the most tickets."
- Health Distribution pie: "Critical and At Risk segments are where Support and CS coordinate most."

For TAE Director:
- ARR Waterfall: "TAEs play a critical role in the expansion number."
- Renewal Rate Trend: "If this dips below target, TAE engagement strategies may need adjustment."

TRANSITION: "Now let's get tactical. How do we forecast which renewals will close?"
""")


# ============================================================
# SLIDE 6: LIVE DEMO — Renewal Forecast
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)
add_shape_bg(slide, Inches(0), Inches(3.5), Inches(13.333), Inches(0.06), GREEN)

add_text_box(slide, Inches(1), Inches(0.5), Inches(3), Inches(0.5),
             "LIVE DEMO", 14, GREEN, True)
add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1),
             "Renewal Forecast & Pipeline", 44, WHITE, True)
add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.5),
             "Forecast renewal outcomes by quarter, segment, and category", 20, RGBColor(148, 163, 184))

add_bullet_points(slide, Inches(1), Inches(4.0), Inches(11), Inches(3), [
    "Filterable by Segment, Region, CSM — self-service for all teams",
    "Commit vs At Risk gap = the CS Ops priority",
    "Monthly timeline prevents last-minute fire drills",
    "Exportable data — no more waiting for weekly Ops emails",
], 16, RGBColor(203, 213, 225))

add_notes(slide, """SWITCH TO DASHBOARD — Renewal Forecast page

Show filters: "Everything filterable by segment, region, CSM. CCO can slice for board reporting. TAEs can see their portfolio. Support can correlate with workload by region."

Metrics row: "This is our pipeline. [X] renewals representing $[X]M in ARR."

Forecast by Category: "This is how I'd run the weekly renewal review. Start with category view, then drill into At Risk."

Renewal Timeline: "Month-by-month shows workload distribution. Heavy months need front-loaded TAE engagement."

Data table + download: "Everything is exportable. CSMs and TAEs pull their own lists."

TRANSITION: "But what about accounts that don't renew? Let's look at what we learn from churn."
""")


# ============================================================
# SLIDE 7: LIVE DEMO — Churn Analysis
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)
add_shape_bg(slide, Inches(0), Inches(3.5), Inches(13.333), Inches(0.06), RED)

add_text_box(slide, Inches(1), Inches(0.5), Inches(3), Inches(0.5),
             "LIVE DEMO", 14, RED, True)
add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1),
             "Churn Analysis & Risk Intelligence", 44, WHITE, True)
add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.5),
             "Why are we losing customers, and can we predict it?", 20, RGBColor(148, 163, 184))

add_bullet_points(slide, Inches(1), Inches(4.0), Inches(5.5), Inches(3), [
    "CCO: Churn reasons drive different strategies — budget vs competitor",
    "Support: Low utilization + high churn = escalation pipeline",
    "TAE: Segment churn informs coverage model decisions",
], 16, RGBColor(203, 213, 225), bold_prefix=True)

add_text_box(slide, Inches(7.5), Inches(4.2), Inches(5), Inches(2),
             'High-Risk Accounts table = early warning system for accounts with >30% churn probability',
             16, RGBColor(148, 163, 184))

add_notes(slide, """SWITCH TO DASHBOARD — Churn Analysis page

For CCO:
- Churn Reasons: "Gold for executive strategy. Budget Constraints = value-proof problem. Competitor Switch = product-market problem. CS Ops surfaces this distinction."
- Logo Churn Rate: "Losing 10 SMB accounts at $10K is different from 1 Enterprise at $100K."

For Support Ops:
- Risk Signals scatter: "Low utilization + high churn probability = escalations happening now or about to."
- High-Risk table: "If your team sees escalations from these names, that's not a coincidence — it's a pattern."

For TAE:
- Churn by Segment: "Higher Mid-Market churn tells us about TAE coverage models."
- Monthly Trend: "TAEs on the ground often have the earliest signal."

TRANSITION: "Churn analysis is backward-looking. Let's flip to the proactive side."
""")


# ============================================================
# SLIDE 8: LIVE DEMO — Account Health
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)
add_shape_bg(slide, Inches(0), Inches(3.5), Inches(13.333), Inches(0.06), CYAN)

add_text_box(slide, Inches(1), Inches(0.5), Inches(3), Inches(0.5),
             "LIVE DEMO", 14, CYAN, True)
add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1),
             "Account Health Overview", 44, WHITE, True)
add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.5),
             "Health scores are only useful if they drive action", 20, RGBColor(148, 163, 184))

add_bullet_points(slide, Inches(1), Inches(4.0), Inches(11), Inches(3), [
    "Healthy Accounts % = North Star metric for CS Ops",
    "Utilization vs NPS scatter reveals different playbooks:",
    "    High utilization + Low NPS = Support problem (experience)",
    "    Low utilization + High NPS = Adoption problem (enablement)",
    "CSM Portfolio Health table — identify rebalancing opportunities",
], 16, RGBColor(203, 213, 225))

add_notes(slide, """SWITCH TO DASHBOARD — Account Health page

Key callouts:
- Healthy Accounts %: "Our North Star. Every initiative should move this number up."
- Utilization vs NPS: "One of the most powerful views. Different quadrants = different playbooks."

For TAE Director — CSM Portfolio Health table:
"Critical for TAE and CSM leadership. If one CSM has 40% at-risk, that's a training or rebalancing conversation."

For CCO:
"Health Score Trend shows if our book is getting healthier or sicker. If flat while growing ARR, we have a ticking time bomb."

TRANSITION: "Now let me show you something I'm particularly proud of — the EBR workflow."
""")


# ============================================================
# SLIDE 9: LIVE DEMO — EBR Preparation (THE WOW MOMENT)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)
add_shape_bg(slide, Inches(0), Inches(3.5), Inches(13.333), Inches(0.06), PURPLE)

add_text_box(slide, Inches(1), Inches(0.3), Inches(4), Inches(0.5),
             "LIVE DEMO — THE WOW MOMENT", 14, PURPLE, True)
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(1),
             "EBR Preparation Workflow", 44, WHITE, True)
add_text_box(slide, Inches(1), Inches(2.3), Inches(11), Inches(0.5),
             "Turns a 2-hour prep process into 2 minutes", 22, RGBColor(148, 163, 184))

add_bullet_points(slide, Inches(1), Inches(3.5), Inches(5.5), Inches(3.5), [
    "Auto-generated EBR brief with 5 sections:",
    "   1. Account Overview — ARR, segment, health, NPS",
    "   2. Engagement & Adoption — utilization gauge, support data",
    "   3. Risk Assessment — auto-flagged risks + action items",
    "   4. Expansion Opportunities — cross-sell, PS, multi-year",
    "   5. Executive Talking Points — opening to the ask",
], 16, RGBColor(203, 213, 225))

add_text_box(slide, Inches(7.5), Inches(3.8), Inches(5), Inches(3),
             "DEMO STEPS:\n1. Select a large Enterprise account\n2. Click \"Generate EBR Brief\"\n3. Let animation run\n4. Walk through each section\n5. Show download button\n\nPAUSE for reactions here.",
             15, AMBER)

add_notes(slide, """*** THIS IS YOUR STRONGEST MOMENT — SPEND TIME HERE ***

SWITCH TO DASHBOARD — EBR Preparation page
1. Select a large Enterprise account
2. Click "Generate EBR Brief"
3. Let the progress animation run (builds anticipation)

Walk through each section:

Account Overview: "Everything a CSM or TAE needs before walking into the room."

Engagement: "[TAE Director], this gauge is what your TAEs see. Below 40% = adoption conversation, not renewal."

Risk Assessment: "CS Ops adds the most value here. System auto-flags risks and generates action items. The playbook is built in."

For Support Ops: "Notice it pulls open tickets and escalation data. Support activity is a first-class signal, not an afterthought."

Expansion: "This ties CS to revenue. Cross-sell, multi-year, PS upsell — all surfaced automatically."

Talking Points: "Actual talking points for the meeting. Opening, value delivered, strategic alignment, the ask."

Download: "One click to export. Share with TAE, CSM, exec sponsor — everyone walks in prepared."

TRANSITION: "Let me show PS & Training, then we'll close with automated alerts."
""")


# ============================================================
# SLIDE 10: LIVE DEMO — PS & Training
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)
add_shape_bg(slide, Inches(0), Inches(3.5), Inches(13.333), Inches(0.06), AMBER)

add_text_box(slide, Inches(1), Inches(0.5), Inches(3), Inches(0.5),
             "LIVE DEMO", 14, AMBER, True)
add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1),
             "PS & Training Analytics", 44, WHITE, True)

add_bullet_points(slide, Inches(1), Inches(3.0), Inches(11), Inches(4), [
    "TAE: PS utilization by account — unused hours = risk, depleted = upsell",
    "CCO: PS Revenue trend feeds into gross margin and time-to-value",
    "Training Impact on Health: trained accounts show higher NPS,",
    "    higher utilization, and lower churn probability",
    "This is the data-driven case for training investment",
], 16, RGBColor(203, 213, 225), bold_prefix=True)

add_notes(slide, """SWITCH TO DASHBOARD — PS & Training page

For TAE Director:
"PS is often owned or influenced by TAE teams. Unused PS hours = risk of feeling they didn't get value. 80%+ used = upsell opportunity."

Training Impact table:
"This is the ROI case for training investment. Trained accounts have higher NPS, higher utilization, lower churn probability."

For CCO:
"PS Revenue trend shows whether we're capturing services revenue effectively."

TRANSITION: "And finally — the piece that ties everything together."
""")


# ============================================================
# SLIDE 11: LIVE DEMO — Automated Alerts
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)
add_shape_bg(slide, Inches(0), Inches(3.5), Inches(13.333), Inches(0.06), RED)

add_text_box(slide, Inches(1), Inches(0.5), Inches(3), Inches(0.5),
             "LIVE DEMO", 14, RED, True)
add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1),
             "Automated Alerts Engine", 44, WHITE, True)
add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.5),
             "From reporting to intelligence — replacing manual weekly reviews", 20, RGBColor(148, 163, 184))

# Alert tiers
add_shape_bg(slide, Inches(1), Inches(3.8), Inches(3.5), Inches(1.2), RGBColor(50, 20, 30))
add_text_box(slide, Inches(1.3), Inches(3.9), Inches(3), Inches(0.4), "CRITICAL", 14, RED, True)
add_text_box(slide, Inches(1.3), Inches(4.3), Inches(3), Inches(0.5),
             "Renewal <60 days + churn >20%", 13, RGBColor(203, 213, 225))

add_shape_bg(slide, Inches(5), Inches(3.8), Inches(3.5), Inches(1.2), RGBColor(40, 30, 10))
add_text_box(slide, Inches(5.3), Inches(3.9), Inches(3), Inches(0.4), "HIGH", 14, AMBER, True)
add_text_box(slide, Inches(5.3), Inches(4.3), Inches(3), Inches(0.5),
             "Critical health OR 2+ escalations", 13, RGBColor(203, 213, 225))

add_shape_bg(slide, Inches(9), Inches(3.8), Inches(3.5), Inches(1.2), RGBColor(10, 30, 45))
add_text_box(slide, Inches(9.3), Inches(3.9), Inches(3), Inches(0.4), "MEDIUM", 14, CYAN, True)
add_text_box(slide, Inches(9.3), Inches(4.3), Inches(3), Inches(0.5),
             "No login for 30+ days", 13, RGBColor(203, 213, 225))

add_bullet_points(slide, Inches(1), Inches(5.5), Inches(11), Inches(1.5), [
    "Support Ops: Alert distribution by CSM reveals portfolio imbalances",
    "CCO: Weekly Digest = Monday morning leadership email",
    "TAE: Filter by accounts for data-driven weekly prioritization",
], 16, RGBColor(203, 213, 225), bold_prefix=True)

add_notes(slide, """SWITCH TO DASHBOARD — Automated Alerts page

Explain the three tiers:
"Critical: Renewal within 60 days AND churn probability above 20%. Executive intervention now.
High: Critical health score OR 2+ escalations. Health check within 48 hours.
Medium: No login for 30+ days. Re-engagement needed."

For Support Ops: "Alert Distribution by CSM chart. Disproportionate alerts = rebalancing signal. Connect dots between Support escalations and CS alerts."

For CCO — Weekly Digest: "This is what I'd send every Monday. Total alerts, ARR impacted, top actions. No one needs to log in for the headline."

For TAE: "Filter by their accounts. Data-driven prioritization, no guessing."

TRANSITION: "Let me close with how I see this fitting into the broader vision."
""")


# ============================================================
# SLIDE 12: Vision & Impact (Closing)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.8), NAVY)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(1),
             "Vision & Impact", 40, WHITE, True)

# Four pillars
pillars = [
    ("1", "Data Drives Action", "Every page ends with something\nsomeone can DO — export, generate,\nact on an alert.", CYAN),
    ("2", "Cross-Functional\nVisibility", "Support, TAE, and executive data\nwoven into every view.\nNo silos.", GREEN),
    ("3", "Automation Replaces\nManual Toil", "Alerts save hours of weekly review.\nEBR generator cuts prep by 80%.\nCSMs spend time with customers.", PURPLE),
    ("4", "Built to Scale", "250 accounts today.\nSame framework works with\n2,500 or 25,000.", AMBER),
]

for i, (num, title, desc, color) in enumerate(pillars):
    x = Inches(0.6 + i * 3.2)
    add_shape_bg(slide, x, Inches(2.3), Inches(2.9), Inches(4.5), LIGHT_BG)
    add_shape_bg(slide, x, Inches(2.3), Inches(2.9), Inches(0.08), color)

    add_text_box(slide, x + Inches(0.3), Inches(2.6), Inches(2.3), Inches(0.5),
                 num, 28, color, True)
    add_text_box(slide, x + Inches(0.3), Inches(3.2), Inches(2.3), Inches(0.8),
                 title, 17, NAVY, True)
    add_text_box(slide, x + Inches(0.3), Inches(4.2), Inches(2.3), Inches(2),
                 desc, 13, MUTED_TEXT)

add_notes(slide, """CLOSING SCRIPT — deliver with confidence, not rushed:

"What I've shown you today is a prototype — but it represents how I think about CS Operations:

1. Data should drive action, not just reporting. Every page ends with something someone can do.

2. Cross-functional visibility is essential. [Support Director], your team's data is woven into every view. [TAE Director], TAE portfolio health is front and center. [CCO], the executive layer gives you the board-ready view.

3. Automation replaces manual toil. The alerts engine saves hours. The EBR generator cuts prep by 80%. That's CS Ops — freeing up CSMs and TAEs for customers, not spreadsheets.

4. Scalability. This was built with 250 accounts. The same framework works with 25,000. As Sumo Logic grows, operational infrastructure scales with it. That's the role I want to play.

I built this because the best way to show what I can do is to actually do it. I'm excited about bringing this operational thinking to Sumo Logic's Customer Success organization.

I'd love to hear your questions."
""")


# ============================================================
# SLIDE 13: First 90 Days
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), CYAN)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.7),
             "IF YOU ASK...", 12, CYAN, True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.8),
             "My First 90 Days", 32, NAVY, True)

# Timeline
phases = [
    ("Weeks 1-2", "DISCOVER", "Audit current data sources,\nreporting cadence, and tools\n(Gainsight, Salesforce, etc.)", CYAN),
    ("Weeks 3-4", "IDENTIFY", "Find the #1 manual process\nto automate. Map stakeholder\nneeds across CS, Support, TAE.", GREEN),
    ("Month 2", "BUILD", "Deliver the first operational\ndashboard with real Sumo Logic\ndata. Quick wins first.", PURPLE),
    ("Month 3", "LAUNCH", "Deploy automated alert system.\nMeasure time saved.\nIterate based on feedback.", AMBER),
]

for i, (time, phase, desc, color) in enumerate(phases):
    x = Inches(0.8 + i * 3.1)
    add_shape_bg(slide, x, Inches(2.2), Inches(2.8), Inches(4.2), LIGHT_BG)
    add_shape_bg(slide, x, Inches(2.2), Inches(2.8), Inches(0.08), color)
    add_text_box(slide, x + Inches(0.2), Inches(2.5), Inches(2.4), Inches(0.4),
                 time, 13, color, True)
    add_text_box(slide, x + Inches(0.2), Inches(3.0), Inches(2.4), Inches(0.5),
                 phase, 22, NAVY, True)
    add_text_box(slide, x + Inches(0.2), Inches(3.8), Inches(2.4), Inches(2.5),
                 desc, 14, MUTED_TEXT)

add_notes(slide, """Only show this slide if asked "What would you do in your first 90 days?"

Otherwise, have it ready as a backup.

Key message: You move fast, start with discovery, and deliver tangible value by month 2-3.
""")


# ============================================================
# SLIDE 14: Q&A Preparation (hidden — for your reference)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)

add_text_box(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.7),
             "Q&A PREPARATION — FOR YOUR EYES ONLY", 14, RED, True)
add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
             "Anticipated Questions & Answers", 32, WHITE, True)

add_bullet_points(slide, Inches(1), Inches(2.5), Inches(11), Inches(4.5), [
    'CCO: "How would you prioritize features?" → Alerts first (fastest ROI), then Renewal Forecast, then EBR Prep.',
    'CCO: "How does this connect to existing tools?" → Proof of concept. Production pulls from Gainsight, Salesforce, Zendesk, product telemetry.',
    'CCO: "What\'s missing?" → Product usage telemetry, predictive ML for churn, Slack/email alert delivery.',
    '',
    'Support: "How to integrate support data deeper?" → Ticket severity, resolution time, CSAT per ticket, first-response time.',
    'Support: "CS Ops + Support Ops collaboration?" → Shared visibility, joint playbooks for at-risk accounts.',
    '',
    'TAE: "How would TAEs use this daily?" → Renewal Forecast filtered to their accounts + EBR Prep in 2 min vs 2 hours.',
    'TAE: "How to measure TAE impact?" → Expansion ARR for TAE-engaged vs non-engaged accounts, PS utilization, renewal rates.',
    '',
    'General: "Tech stack?" → Python, Streamlit, Plotly, Pandas, Docker. In production, use whatever the org supports.',
], 14, RGBColor(203, 213, 225), bold_prefix=True)

add_notes(slide, """DO NOT SHOW THIS SLIDE.
This is your cheat sheet for Q&A. Keep it in Presenter View or memorize the key answers.

Additional prep:
- "How long did this take?" → Be honest. Frame as: "I invested [X] because I wanted to demonstrate I can execute, not just talk."
- If they ask about something you don't know about Sumo Logic's internal tools, say: "I'd love to learn more about your current stack — my approach is tool-agnostic and adapts to whatever systems are in place."
""")


# ============================================================
# SLIDE 15: Thank You / Close
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, NAVY)
add_shape_bg(slide, Inches(0), Inches(6.8), Inches(13.333), Inches(0.06), CYAN)

add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
             "Thank You", 52, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.8),
             "I built this because the best way to show\nwhat I can do is to actually do it.", 22, RGBColor(148, 163, 184), alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
             "Questions?", 28, CYAN, True, PP_ALIGN.CENTER)

add_notes(slide, """End with energy and confidence. Make eye contact. Smile.

You built this — you know it better than anyone in that room. Good luck!
""")


# ============================================================
# SAVE
# ============================================================
output_path = "CS_Ops_Panel_Presentation.pptx"
prs.save(output_path)
print(f"\nPresentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
print("\nSlide overview:")
slide_names = [
    "1.  Title Slide",
    "2.  Agenda",
    "3.  Your Panel — Tailored Value (reference only)",
    "4.  Why I Built This Dashboard",
    "5.  LIVE DEMO: Executive Summary",
    "6.  LIVE DEMO: Renewal Forecast",
    "7.  LIVE DEMO: Churn Analysis",
    "8.  LIVE DEMO: Account Health",
    "9.  LIVE DEMO: EBR Preparation (WOW moment)",
    "10. LIVE DEMO: PS & Training",
    "11. LIVE DEMO: Automated Alerts",
    "12. Vision & Impact (Closing)",
    "13. First 90 Days (backup slide)",
    "14. Q&A Cheat Sheet (DO NOT SHOW)",
    "15. Thank You",
]
for name in slide_names:
    print(f"  {name}")
print("\nTIP: All speaker notes are in Presenter View (View > Notes in PowerPoint)")
